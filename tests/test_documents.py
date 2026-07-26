"""Tests for read_document: PDF/docx/xlsx/pptx/csv text extraction.

docx and xlsx are built as minimal OOXML zips with the stdlib (the reader
parses the same XML members it would in a real file), pptx/csv use the
libraries actually installed, and the PDF path is checked for its graceful
"install a library" message.
"""

import importlib.util
import tempfile
import unittest
import zipfile
from pathlib import Path

from jarvis.tools import documents

_DOCX_XML = (
    '<?xml version="1.0"?>'
    '<w:document xmlns:w="http://schemas.openxmlformats.org/'
    'wordprocessingml/2006/main"><w:body>'
    '<w:p><w:r><w:t>Hello</w:t></w:r><w:r><w:t> World</w:t></w:r></w:p>'
    '<w:p><w:r><w:t>Second line</w:t></w:r></w:p>'
    '</w:body></w:document>'
)

_SST_XML = (
    '<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
    '<si><t>Name</t></si><si><t>Alice</t></si></sst>'
)
_SHEET_XML = (
    '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/'
    '2006/main"><sheetData>'
    '<row><c r="A1" t="s"><v>0</v></c><c r="B1"><v>42</v></c></row>'
    '<row><c r="A2" t="s"><v>1</v></c>'
    '<c r="B2" t="inlineStr"><is><t>hi</t></is></c></row>'
    '</sheetData></worksheet>'
)


class ReadDocumentTests(unittest.TestCase):
    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())

    def tearDown(self):
        import shutil
        shutil.rmtree(self.tmp, ignore_errors=True)

    # -- dispatch guards -------------------------------------------------- #
    def test_missing_file(self):
        self.assertIn("file not found",
                      documents.read_document(str(self.tmp / "nope.pdf")))

    def test_legacy_binary_is_redirected(self):
        legacy = self.tmp / "old.doc"
        legacy.write_bytes(b"\xd0\xcf\x11\xe0")   # OLE magic; content irrelevant
        out = documents.read_document(str(legacy))
        self.assertIn("legacy binary", out)
        self.assertIn(".docx", out)

    def test_plain_text_falls_through_to_read_file(self):
        txt = self.tmp / "notes.txt"
        txt.write_text("just some text", encoding="utf-8")
        self.assertIn("just some text", documents.read_document(str(txt)))

    # -- docx (stdlib zip+xml) ------------------------------------------- #
    def test_docx_paragraphs_and_runs(self):
        docx = self.tmp / "doc.docx"
        with zipfile.ZipFile(docx, "w") as z:
            z.writestr("word/document.xml", _DOCX_XML)
        out = documents.read_document(str(docx))
        self.assertIn("Hello World", out)      # two runs joined in one paragraph
        self.assertIn("Second line", out)

    # -- xlsx (stdlib zip+xml) ------------------------------------------- #
    def test_xlsx_shared_inline_and_numeric_cells(self):
        xlsx = self.tmp / "book.xlsx"
        with zipfile.ZipFile(xlsx, "w") as z:
            z.writestr("xl/sharedStrings.xml", _SST_XML)
            z.writestr("xl/worksheets/sheet1.xml", _SHEET_XML)
        out = documents.read_document(str(xlsx))
        self.assertIn("sheet: sheet1", out)
        self.assertIn("Name", out)   # shared string, index 0
        self.assertIn("Alice", out)  # shared string, index 1
        self.assertIn("42", out)     # bare numeric cell
        self.assertIn("hi", out)     # inline string

    # -- pdf graceful degradation ---------------------------------------- #
    @unittest.skipIf(importlib.util.find_spec("pypdf")
                     or importlib.util.find_spec("PyPDF2"),
                     "a PDF library is installed; the missing-lib path is moot")
    def test_pdf_without_library_explains_how_to_install(self):
        pdf = self.tmp / "x.pdf"
        pdf.write_bytes(b"%PDF-1.4 fake")
        out = documents.read_document(str(pdf))
        self.assertIn("pip install pypdf", out)

    # -- pptx (python-pptx is installed) --------------------------------- #
    @unittest.skipUnless(importlib.util.find_spec("pptx"),
                         "python-pptx not installed")
    def test_pptx_slide_text(self):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "Quarterly results"
        pptx_path = self.tmp / "deck.pptx"
        prs.save(str(pptx_path))
        out = documents.read_document(str(pptx_path))
        self.assertIn("slide 1", out)
        self.assertIn("Quarterly results", out)

    # -- csv (pandas is installed) --------------------------------------- #
    @unittest.skipUnless(importlib.util.find_spec("pandas"),
                         "pandas not installed")
    def test_csv_preview_reports_shape_and_values(self):
        csv = self.tmp / "data.csv"
        csv.write_text("name,age\nAlice,30\nBob,25\n", encoding="utf-8")
        out = documents.read_document(str(csv))
        self.assertIn("2 columns", out)
        self.assertIn("name", out)
        self.assertIn("Alice", out)


class ExposureTests(unittest.TestCase):
    def test_wired_into_registry_and_agents(self):
        from jarvis.tools.registry import _HANDLERS
        from jarvis.tools.schema import ACTIONS_BY_NAME
        from jarvis.agent.coder import ALLOWED
        from jarvis.agent.subagent import HEADLESS
        self.assertIn("read_document", ACTIONS_BY_NAME)
        self.assertIn("read_document", _HANDLERS)
        self.assertIn("read_document", ALLOWED)
        self.assertIn("read_document", HEADLESS)


if __name__ == "__main__":
    unittest.main()
