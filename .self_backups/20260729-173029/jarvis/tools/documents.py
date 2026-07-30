"""Read the text out of common document formats.

``read_document(path)`` dispatches by extension so tasks like "summarise this
PDF" or "what's in that spreadsheet" work headlessly. Office formats (docx,
xlsx) are just zipped XML, so they are parsed with the stdlib - no dependency
needed. PDF needs a third-party library; if none is installed the reader says
which one to ``pip install`` instead of failing cryptically.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from .files import _expand, read_file

# OOXML namespaces.
_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
_S = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"

# Old binary Office formats need a different toolchain than their XML successors.
_LEGACY = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}


class _MissingLib(Exception):
    """Raised when a format needs a library that is not installed; the message
    tells the user exactly what to pip-install."""


def read_document(path: str, max_chars: int = 20000) -> str:
    p = _expand(path)
    if not p.exists():
        return f"file not found: {p}"
    if not p.is_file():
        return f"not a file: {p}"
    try:
        max_chars = max(500, min(200000, int(max_chars or 20000)))
    except (TypeError, ValueError):
        max_chars = 20000

    ext = p.suffix.lower()
    if ext in _LEGACY:
        return (f"{p.name} is a legacy binary {ext} file - save it as "
                f"{_LEGACY[ext]} and read that instead.")
    try:
        if ext == ".pdf":
            text = _read_pdf(p)
        elif ext == ".docx":
            text = _read_docx(p)
        elif ext in (".xlsx", ".xlsm"):
            text = _read_xlsx(p)
        elif ext == ".pptx":
            text = _read_pptx(p)
        elif ext in (".csv", ".tsv"):
            text = _read_csv(p)
        else:
            return read_file(str(p), max_chars=max_chars)   # plain text
    except _MissingLib as exc:
        return str(exc)
    except Exception as exc:
        return f"could not read {p.name}: {exc}"

    text = (text or "").strip()
    if not text:
        return (f"{p.name} has no extractable text (it may be scanned images - "
                f"that would need OCR).")
    if len(text) > max_chars:
        text = text[:max_chars] + f"\n...[truncated, {len(text)} chars total]"
    return f"content of {p.name}:\n{text}"


def _read_pdf(p: Path) -> str:
    """Extract text with whichever PDF library is installed."""
    reader_cls = None
    for mod in ("pypdf", "PyPDF2"):
        try:
            reader_cls = __import__(mod, fromlist=["PdfReader"]).PdfReader
            break
        except Exception:
            continue
    if reader_cls is None:
        raise _MissingLib("reading PDFs needs a PDF library - install one "
                          "with:  pip install pypdf")
    reader = reader_cls(str(p))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            t = (page.extract_text() or "").strip()
        except Exception:
            t = ""
        if t:
            parts.append(f"--- page {i} ---\n{t}")
    return "\n\n".join(parts)


def _read_docx(p: Path) -> str:
    """docx text lives in word/document.xml as <w:t> runs grouped by <w:p>
    paragraphs - parse it straight from the zip, no python-docx needed."""
    import xml.etree.ElementTree as ET

    with zipfile.ZipFile(p) as z:
        root = ET.fromstring(z.read("word/document.xml"))
    lines = ["".join(t.text or "" for t in para.iter(f"{_W}t"))
             for para in root.iter(f"{_W}p")]
    return "\n".join(lines)


def _read_xlsx(p: Path) -> str:
    """Parse the xlsx OOXML zip directly: shared strings + each sheet's cells.

    ponytail: presents cells in document order joined by tabs - good enough to
    read a sheet. Add openpyxl if exact column/gap fidelity ever matters.
    """
    import xml.etree.ElementTree as ET

    with zipfile.ZipFile(p) as z:
        names = z.namelist()
        shared: list[str] = []
        if "xl/sharedStrings.xml" in names:
            sst = ET.fromstring(z.read("xl/sharedStrings.xml"))
            shared = ["".join(t.text or "" for t in si.iter(f"{_S}t"))
                      for si in sst.iter(f"{_S}si")]
        out = []
        sheets = sorted(n for n in names
                        if n.startswith("xl/worksheets/") and n.endswith(".xml"))
        for sheet in sheets:
            ws = ET.fromstring(z.read(sheet))
            rows = []
            for row in ws.iter(f"{_S}row"):
                cells = [_cell_value(c, shared) for c in row.iter(f"{_S}c")]
                if any(v.strip() for v in cells):
                    rows.append("\t".join(cells))
            if rows:
                out.append(f"--- sheet: {Path(sheet).stem} ---\n" + "\n".join(rows))
        return "\n\n".join(out)


def _cell_value(c, shared: list[str]) -> str:
    import xml.etree.ElementTree as ET  # noqa: F401  (Element type is duck-typed)

    t = c.get("t")
    v = c.find(f"{_S}v")
    if t == "s" and v is not None and v.text is not None:      # shared string
        try:
            return shared[int(v.text)]
        except (ValueError, IndexError):
            return ""
    if t == "inlineStr":                                       # inline string
        is_ = c.find(f"{_S}is")
        return "".join(x.text or "" for x in is_.iter(f"{_S}t")) if is_ is not None else ""
    return v.text if v is not None and v.text else ""           # number/bool/date


def _read_pptx(p: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        raise _MissingLib("reading PowerPoint needs python-pptx - install it "
                          "with:  pip install python-pptx")
    prs = Presentation(str(p))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text.strip() for shape in slide.shapes
                 if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        if texts:
            out.append(f"--- slide {i} ---\n" + "\n".join(texts))
    return "\n\n".join(out)


def _read_csv(p: Path) -> str:
    """A structured preview (shape + columns + head) via pandas when it is
    installed; plain text otherwise (a CSV is readable either way)."""
    try:
        import pandas as pd  # type: ignore
    except Exception:
        return read_file(str(p))
    sep = "\t" if p.suffix.lower() == ".tsv" else ","
    df = pd.read_csv(p, sep=sep, nrows=1000, dtype=str, keep_default_na=False)
    head = df.head(20).to_string(index=False)
    return (f"{df.shape[0]} rows read (up to 1000) x {df.shape[1]} columns: "
            f"{list(df.columns)}\n\n{head}")
